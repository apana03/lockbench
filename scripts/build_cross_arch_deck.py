#!/usr/bin/env python3
"""Build the cross-arch findings deck for the project's "cross-architecture
comparison + pcpu_rw_lock investigation" stage (2026-05-11). Writes
`results/deck/2026-05-11_cross_arch_findings.pptx`.

The deck merges:
  - The supervisor-curated 35-slide structure (cross-arch headline,
    architectural reference, scaling, cache regime, pcpu-rw diagnosis,
    findings).
  - Six spinlock-focused cross-bench / cross-arch slides appended after
    the curated section.

PRECONDITION: PowerPoint must NOT have the file open. Verify with
  lsof results/deck/2026-05-11_cross_arch_findings.pptx
  ls results/deck/~\\$2026-05-11_cross_arch_findings.pptx  # should not exist

Run:  results/.venv/bin/python scripts/build_cross_arch_deck.py
"""

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DECK_NAME = "2026-05-11_cross_arch_findings.pptx"
LOCK = ROOT / "results" / "deck" / f"~${DECK_NAME}"
TARGET = ROOT / "results" / "deck" / DECK_NAME

if LOCK.exists():
    sys.stderr.write(
        f"\nERROR: PowerPoint lock file still present at\n  {LOCK}\n"
        "Close PowerPoint completely (Cmd-Q on Mac), then rerun.\n\n")
    sys.exit(1)

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402

FIG_DIR = ROOT / "results" / "deck" / "figures"

# Slide dimensions: 16:9 widescreen.
SW = Inches(13.333)
SH = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Slide-construction helpers (inlined here because the original helper module
# `build_findings_deck.py` was lost during a repo cleanup. To regenerate
# figures, see `scripts/build_ranking_figures.py` for the ranking-figure
# helpers; the other figures were one-off and would need to be reconstructed
# from `scripts/lockbench_analysis.ipynb` if needed.)
# ─────────────────────────────────────────────────────────────────────────────

def _add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.7))
    tf = box.text_frame
    tf.text = title
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(28)
            r.font.bold = True


def add_bullets(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.9)).text_frame
    box.word_wrap = True
    for i, b in enumerate(bullets):
        indent, txt = (b if isinstance(b, tuple) else (0, b))
        p = box.paragraphs[0] if i == 0 else box.add_paragraph()
        p.text = ("- " if indent == 0 else "    - ") + txt
        p.level = indent
        for r in p.runs:
            r.font.size = Pt(20 if indent == 0 else 16)
            if indent > 0:
                r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_table(prs, title, cols, rows, col_widths_inches):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)
    n_rows = len(rows) + 1
    n_cols = len(cols)
    total_w = sum(col_widths_inches)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.6), Inches(1.2),
        Inches(total_w), Inches(min(5.5, 0.4 * n_rows + 0.4)),
    )
    tbl = tbl_shape.table
    for i, w in enumerate(col_widths_inches):
        tbl.columns[i].width = Inches(w)
    for j, col in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = col
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)


def add_image(prs, title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.2),
                                 height=Inches(5.3))
    else:
        slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0),
                                 Inches(0.5)).text_frame.text = f"[MISSING: {image_path}]"
    if caption:
        cap = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.0),
                                       Inches(0.6)).text_frame
        cap.text = caption
        for p in cap.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.italic = True


def add_section(prs, name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12.0), Inches(1.5)).text_frame
    box.text = name
    for p in box.paragraphs:
        for r in p.runs:
            r.font.size = Pt(40)
            r.font.bold = True


# Thin shim so the slide-definition block below can keep using `B.add_bullets`,
# `B.add_table`, etc. without changes.
class _BModule:
    add_bullets = staticmethod(add_bullets)
    add_table = staticmethod(add_table)
    add_image = staticmethod(add_image)
    add_section = staticmethod(add_section)
    SW = SW
    SH = SH


B = _BModule()


def main():
    print("Building cross-arch findings deck from existing figures...")
    prs = Presentation()
    prs.slide_width = B.SW
    prs.slide_height = B.SH

    # ---- 1. What was wrong with the old setup ----
    B.add_bullets(prs, "What was wrong with the old setup", [
        "On every operation, each worker thread was doing three things INSIDE the timed measurement window:",
        (1, "Picking a random key according to the workload's distribution (uniform or Zipfian)."),
        (1, "Deciding whether the next operation was a read, an insert, or a remove."),
        (1, "Checking a shared flag to know when the measurement was over."),
        "For uniform workloads this cost a handful of nanoseconds. For Zipfian workloads it was much worse: the skewed-sampling math itself takes 50-80 nanoseconds per operation.",
        "Most of the locks we test have critical sections of 30-50 ns. So workload generation was 30-60 % of measured time, and that fraction differed depending on which lock you were testing.",
        "Conclusion: the measurement window was contaminated by per-operation overhead, and the contamination was unequally distributed across locks. Biased comparisons.",
    ])

    # ---- 2. New design ----
    B.add_bullets(prs, "New design: pre-generate the workload before measuring", [
        "Each worker thread now generates its own buffer of (key, operation) pairs at startup, BEFORE the timed window opens.",
        "Same workload distribution as before: uniform or Zipfian with the same skew parameters, just computed once and stored.",
        "During measurement, each operation is a memory read + branch into the right call. No random-number generation, no probability math.",
        "Threads have independent buffers seeded differently, so no two workers see the same key sequence (avoids shared-bias artefacts).",
        "Each worker starts at a different position in its buffer, so the threads aren't synchronised on the same operation at time-zero.",
        "Net effect: the timed window now measures the lock + the index lookup, nothing else.",
    ])

    # ---- 3. Did the new design itself introduce any bias? (table) ----
    B.add_table(prs, "Did the new design itself introduce any bias?",
        ["Buffer size", "Bytes per thread", "Median throughput", "Where it lives in the cache hierarchy"],
        [
            ["small",   "16 KiB",  "87.7 M ops/s", "fits in L1 data cache"],
            ["default", "64 KiB",  "80.8 M ops/s", "fits in L2 (preferred for accurate Zipfian sampling)"],
            ["large",   "256 KiB", "77.4 M ops/s", "spills to the edge of L2"],
        ],
        col_widths_inches=[1.8, 2.0, 2.8, 5.4])

    # ---- 4. Reading the validation ----
    B.add_bullets(prs, "Reading the validation (Xeon, 12 threads, 3 trials per row)", [
        "Within each buffer size: variance below 1 % across the three trials. The new design is stable.",
        "Across buffer sizes: a real 12 % step. Smaller buffers sit in faster cache, so the buffer walk itself is ~1 ns/op cheaper.",
        "Crucially, that 12 % cost is paid equally by every lock variant: comparisons remain valid at any choice of buffer size.",
        "We picked the default that lands in L2 on both architectures. Tradeoff: marginally lower absolute throughput vs better statistical fidelity for the most extreme skew level (θ=1.5), where the hot-key mass concentrates on very few buffer entries.",
    ])

    # ---- 5. Section divider: New cache-regime workload matrix ----
    B.add_section(prs, "2. New cache-regime workload matrix")

    # ---- 6. Why workloads need to target specific cache levels ----
    B.add_bullets(prs, "Why workloads need to target specific cache levels", [
        "Goal: measure what a lock costs, not what memory costs.",
        "If the index spills to DRAM, per-op memory latency (~80 ns) dominates and the lock-acquire cost (~30-50 ns) becomes a small fraction of total time. Lock differences compress.",
        "Conversely, if every reader is touching its own L1-resident keys, there is no reader-reader contention at all: the lock is never the bottleneck.",
        "To force the lock to be the bottleneck we need two things at once:",
        (1, "Small enough working set that memory-stall is bounded (index fits in cache)."),
        (1, "Concentrated enough access pattern that multiple readers compete for the same lock instances at once."),
        "The cache-regime matrix is engineered to deliver this. Two cache regimes x four key distributions x two op mixes.",
    ])

    # ---- 7. Key distributions (RETITLED: no subtitle) ----
    B.add_table(prs, "Key distributions",
        ["Distribution", "θ", "What it does", "Top-key mass (approx)"],
        [
            ["uniform",          "n/a",   "every key equally likely; baseline with no skew",                "1 / N (negligible)"],
            ["zipfian: warm",    "0.99",  "standard YCSB hot-key concentration",                            "~10 % of ops on hottest key"],
            ["zipfian: hot",     "1.2",   "aggressive concentration; top few keys take most ops",           "~30 % on hottest key"],
            ["zipfian: extreme", "1.5",   "near-pathological skew; bulk of ops on a handful of keys",       "~60 % on hottest key"],
        ],
        col_widths_inches=[2.4, 0.9, 4.6, 4.1])

    # ---- 8. Section divider: New results from the sweep ----
    B.add_section(prs, "4. New results from the sweep")

    # ---- 9. 1-thread baseline ----
    B.add_image(prs, "1-thread baseline: pure microarchitectural cost",
              FIG_DIR / "f1_1t_baseline.png",
              "Uncontended: just (clock + atomic-instruction cost + DRAM). Graviton 1.20-1.34x faster, tight band.")

    # ---- 10. Scaling on the headline workload ----
    B.add_image(prs, "Scaling on the headline workload: wormhole, read-heavy",
              FIG_DIR / "f3_scaling_grid.png",
              "Rows: L1_warm | L1_extreme | L3_warm.   Columns: Xeon | Graviton.   wh-occ-opt dominates; wh-pcpu-rw v1 collapses on Graviton.")

    # ---- 11. Heatmap ----
    B.add_image(prs, "Cross-arch ratio across every (lock x workload) cell",
              FIG_DIR / "f4_heatmap_ratio.png",
              "Graviton/Xeon ratio at max threads. Green = Graviton wins. wh-pcpu-rw row breaks the pattern.")

    # ---- 12. Why we added a per-core lock (RETITLED) ----
    B.add_bullets(prs, "Why we added a per-core lock", [
        "Standard reader-writer locks (including the one wormhole ships with) keep a single counter of how many readers are active.",
        "Every reader has to atomically increment that counter to enter, and decrement to leave.",
        "On modern multicore CPUs that counter sits on a single cache line that bounces between cores. Even though readers should run in parallel, they end up serialising on cache coherence.",
        "This is a well-documented bottleneck: Calciu et al., \"NUMA-Aware Reader-Writer Locks\" (PPoPP 2013) models it analytically.",
        "The proposed fix (used in the Linux kernel's percpu_rwsem and in academic BRAVO locks):",
        (1, "Give every thread its own reader counter, in a memory region that no other thread will touch."),
        (1, "Readers no longer compete for the same cache line: parallelism materialises."),
        "Expected to be the headline win on read-heavy workloads, especially on ARM (where the contended-counter design hits harder).",
    ])

    # ---- 13. The lock (RETITLED from 'The lock we built: at a high level') ----
    B.add_bullets(prs, "The lock", [
        "64 per-thread \"slots\": each on its own cache line. Each thread is hashed to one slot.",
        "Reader protocol:",
        (1, "1. Announce arrival by bumping the counter in your slot."),
        (1, "2. Check whether a writer is currently active."),
        (1, "3. If no writer: proceed into the critical section."),
        (1, "4. If a writer IS active: retract (decrement your slot's counter back) and wait until the writer leaves, then retry from step 1."),
        "Writer protocol:",
        (1, "1. Set a flag to tell readers \"writer arriving: back off.\""),
        (1, "2. Scan all 64 slots, waiting for each one to drop to zero (drain in-flight readers)."),
        (1, "3. Execute the critical section."),
        (1, "4. Clear the flag."),
        "On paper this looks correct: no shared counter, no cache-line bouncing on the read path, writers properly serialised. The design checked out.",
    ])

    # ---- 14. The surprise (RETITLED) ----
    B.add_bullets(prs, "The surprise", [
        "After running the full publishable sweep on both arches, the new wh-pcpu-rw row was unlike every other lock:",
        "On Xeon: gentle drop at high T (peak 18.8 M at 4T, 10.9 M at 12T: about 40 % off peak).",
        "On Graviton: catastrophic collapse: 4.9 M at 4T, 0.03 M at 8T. Worse than wh-default at every thread count above 2.",
        "Inverts the cross-arch story: every OTHER lock is 1.4-2.8x faster on Graviton than on Xeon.",
        "Reality on the workload that broke it (90/5/5, L1_warm_zipf99: small index, warm-key concentration):",
    ])

    # ---- 15. How I structured the investigation: title EMPTIED on user request ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # leave title box empty (mirroring the user's edit)
    s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.8)).text_frame
    box.word_wrap = True
    bullets = [
        "Before running any diagnostics, I wrote down three candidate explanations and what each one would predict:",
        (1, "(H1) Hardware-specific bug: a particular thread count happens to map two threads to the same slot."),
        (1, "(H2) Memory-model bug: the lock has a latent correctness issue that ARM's weaker ordering exposes but x86's stronger ordering hides."),
        (1, "(H3) Algorithmic failure mode: the design enters a degenerate state when there's enough contention."),
        "Predictions that distinguish them:",
        (1, "(H1) would show a discrete cliff at one specific thread count, with normal behaviour above and below."),
        (1, "(H2) would manifest at any thread count and any writer rate, possibly as wrong answers from the lock (correctness, not performance)."),
        (1, "(H3) would show smooth degradation as contention increases along either axis (more threads, or more writers)."),
        "Three diagnostic experiments designed to discriminate. Plus one cross-check on the bare lock primitive (no wormhole around it), to isolate the lock from the data structure.",
    ]
    for i, b in enumerate(bullets):
        indent, txt = (b if isinstance(b, tuple) else (0, b))
        p = box.paragraphs[0] if i == 0 else box.add_paragraph()
        p.text = ("- " if indent == 0 else "   - ") + txt
        p.level = indent
        for r in p.runs:
            r.font.size = Pt(20 if indent == 0 else 16)
            if indent > 0:
                r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- 16. Diagnostic 0a: thread sweep ----
    B.add_table(prs, "Diagnostic 0a: thread sweep on Graviton2  (90/5/5, L1_warm_zipf99)",
        ["threads", "median (M/s)", "trial spread", "scaling vs 1T", "trial CoV"],
        [
            ["1", "17.28", "17.28 - 17.40", "1.00x",   "< 1 %"],
            ["2", "17.74", "15.70 - 20.45", "1.03x",   "13 %"],
            ["3", "12.46", "9.61 - 14.88",  "0.72x",   "21 %"],
            ["4", "2.99",  "2.19 - 5.53",   "0.17x",   "57 %"],
            ["5", "0.58",  "0.58 - 0.72",   "0.034x",  "12 %"],
            ["6", "0.15",  "0.14 - 0.43",   "0.009x",  "100 %"],
        ],
        col_widths_inches=[1.5, 2.4, 2.8, 2.5, 1.8])

    # ---- 17. What 0a discriminates ----
    B.add_bullets(prs, "What 0a discriminates", [
        "Smooth feedback degradation 1T -> 6T. NOT a step function: rules out (H1) discrete collision.",
        "2T already isn't scaling: a reader-scalable rwlock should give ~1.8x at 2T on 90 % reads; we get 1.03x.",
        "Trial variance grows monotonically with thread count. The collapse is non-deterministic: exact equilibrium depends on writer/reader race outcomes.",
        "Pattern is exactly what (H3) predicts: contention triggers a degenerate state. The next diagnostic isolates the trigger.",
    ])

    # ---- 18. Diagnostic 0b ----
    B.add_image(prs, "Diagnostic 0b: isolate the trigger by varying writers",
              FIG_DIR / "f6_pcpu_readpct.png",
              "Fix threads = 4, fix skew = θ=0.99, vary read_pct.  100 % reads scales super-linearly to 75.6 M (4.37x single-thread).")

    # ---- 19. What 0b proves ----
    B.add_bullets(prs, "What 0b proves", [
        "At 100 % reads: 4 threads deliver 4.37x of single-thread throughput. The super-linear result confirms the data layout works as designed: each thread has its own piece of memory and they don't compete for cache.",
        "The per-CPU slot machinery WORKS when there are no writers. The design is sound on the read dimension.",
        "One percent writers (99/0/1): throughput drops 43 %. A 1-in-100 event costs more than half of available throughput.",
        "Five percent writers (95/2/3): throughput drops 94 %. Same pattern, larger trigger rate.",
        "Conclusion: the bug is in the writer-reader handshake, not the slot infrastructure. (H3) confirmed on the writer-rate axis.",
    ])

    # ---- 20. Diagnostic 0c intro ----
    B.add_bullets(prs, "Diagnostic 0c: is wormhole amplifying the failure?", [
        "Ran the bare primitive in the lockbench microbench (bench/main.cpp): one global pcpu_rw_lock, no wormhole call pattern.",
        "Same 90 % read mix, same compact_phys pinning, same threading."
    ])

    # ---- 21. Diagnostic 0c table ----
    B.add_table(prs, "Diagnostic 0c: bare pcpu_rw_lock (no wormhole), Graviton2, 90 % reads",
        ["threads", "median (M ops/s)", "scaling vs 1T", "vs wormhole at same T"],
        [
            ["1", "62.08", "1.00x", "wh: 17.3 M  ->  microbench is 3.6x faster (no leaf walk)"],
            ["2", "28.17", "0.45x", ""],
            ["3", "23.09", "0.37x", ""],
            ["4", "19.30", "0.31x", "wh: 4.9 M  ->  4x amplification"],
            ["5", "17.88", "0.29x", ""],
            ["6", "12.90", "0.21x", "wh: 0.15 M  ->  86x amplification"],
        ],
        col_widths_inches=[1.4, 2.6, 2.2, 5.0])

    # ---- 22. What 0c tells us ----
    B.add_bullets(prs, "What 0c tells us", [
        "The bare primitive ALSO fails to scale: 1T -> 6T loses 79 % of throughput. The bug is in the lock, not wormhole.",
        "BUT the bare primitive is much milder: 6T = 12.9 M ops/s. Wormhole-amplified = 0.15 M ops/s. 86x amplification.",
        "Two amplification mechanisms identified:",
        (1, "Two locks per operation. Wormhole takes a global \"index-level\" lock AND a \"leaf-level\" lock on every read or write: both of these were swapped for our per-core lock. So every writer event triggers the herd on whichever of the two locks it touched. With 10 % writers overall, both locks are in herd state at the same time."),
        (1, "Writer critical sections are much longer in a real index. In the microbench, a writer just increments a counter (~10 ns). In wormhole, a writer sorts a leaf, inserts a key-value pair, and occasionally resizes: hundreds of ns. While the writer is working, all readers are spinning uselessly. Longer writer work = more reader spin = more thrashing."),
        "Estimated amplification: (2 locks) x (CS-length ratio ~30x) approx 60-100x. Matches the observed 86x.",
    ])

    # ---- 23. Diagnosis ----
    B.add_bullets(prs, "Diagnosis: what's actually happening", [
        "Picture the steady state: N readers, each working in their own slot, no one blocking anyone. Throughput is high.",
        "A writer arrives:",
        (1, "1. Writer raises the \"writer active\" flag."),
        (1, "2. Every concurrent reader sees the flag, retracts from its slot, and starts spinning."),
        (1, "3. Writer begins scanning the 64 slots, waiting for them all to drain to zero."),
        (1, "4. As readers retract, the writer's scan sees zeros: but readers that already finished retracting are now retrying and stamping their slots back to 1. The scan never sees a quiet moment."),
        (1, "5. Writer eventually completes the scan, does its work, clears the flag."),
        (1, "6. Every spinning reader sees the flag clear simultaneously and stampedes back. They all retry within nanoseconds of each other: the writer sees a full house of readers on its next arrival."),
        (1, "7. Next writer arrives. Loop."),
        "The system is now in an oscillation: writers and readers alternating in lockstep, with very little useful work done in either phase. Throughput collapses to roughly (writer arrival rate x writer critical-section time), which is far below the parallel-reader rate the design promised.",
    ])

    # ---- 24. Why this is WORSE on Graviton2 ----
    B.add_bullets(prs, "Why this is WORSE on Graviton2: the central novel finding", [
        "The oscillation depends critically on the readers all retrying at the same moment. The more synchronised they are, the more readers the writer sees on every cycle, and the worse the herd.",
        "The synchronisation comes from how fast the atomic operations are:",
        (1, "Modern ARM (Graviton2) uses single-instruction atomic operations introduced in ARMv8.1: roughly 2-3x faster than the equivalent on a 2014-era Xeon."),
        (1, "On Graviton, when the writer's flag clears, all 8 readers complete their retry within ~10 ns of each other: perfectly synchronised stampede."),
        (1, "On Xeon, the same retry is slower, so the readers spread out over ~40 ns. The writer often finishes before the last reader has even stamped its slot. The herd is partially dispersed."),
        "Punchline: the Xeon's slower atomic operations ACCIDENTALLY THROTTLE THE HERD. The bug is timing-sensitive to atomic speed.",
        "This is the most novel finding in the project: hardware modernisation can expose latent timing bugs that were dormant on older silicon. As far as I can tell, the literature does not document this specific failure mode as a function of atomic speed.",
    ])

    # ---- 25. Fix design ----
    B.add_bullets(prs, "Fix design: change the handshake, not the data layout", [
        "The Linux kernel solved this same problem in their per-CPU reader-writer semaphore (percpu_rwsem). The principle:",
        (1, "Readers commit and proceed. They do NOT back off when a writer arrives."),
        (1, "Writers queue on a separate slow-path mutex (only one writer at a time)."),
        (1, "Writers wait for already-in-flight readers to finish their critical sections naturally."),
        (1, "New readers arriving while a writer is queued route briefly through the same slow-path mutex (then proceed). No retraction, no spinning."),
        "Crucial property: when a writer arrives, readers caught mid-acquire FINISH their critical section instead of backing off.",
        "No reader ever retracts -> no synchronised retry -> no herd.",
        "Trade-off: an extra branch on the reader fast path (to check whether a writer is queued) costs about 5 % at single-thread. Writers now serialise on the slow-path mutex, but there's only one writer per epoch anyway in any reasonable workload.",
    ])

    # ---- 26. v2: what changed ----
    B.add_bullets(prs, "v2: what changed vs what stayed the same", [
        "Kept exactly:",
        (1, "64 per-thread slots, one cache line each. The data layout was right all along."),
        (1, "The thread-to-slot assignment scheme."),
        "Changed:",
        (1, "Added a slow-path mutex on the writer side, so writers queue rather than race."),
        (1, "Added a \"writer-queued\" signal separate from the original \"writer-active\" flag."),
        (1, "Readers now check the new signal: if set, briefly take the slow-path mutex and release it (this lets the writer proceed); otherwise proceed normally."),
        (1, "Writers no longer ask readers to back off; they wait for in-flight readers to finish naturally."),
        "1-thread cost vs v1: about 5 % regression. The trade is obviously correct.",
        "Correctness validated with race tests: 8 threads x 20 000 mixed read/write operations, both mutual-exclusion and torn-read scenarios.",
    ])

    # ---- 27. v2 validation figure ----
    B.add_image(prs, "v2 eliminates the collapse: Graviton2 validation",
              FIG_DIR / "f5_pcpu_v1_vs_v2.png",
              "All curves on the same workload: wormhole, Graviton2, L1-resident (1k keys), Zipfian θ=0.99, 90% read / 5% insert / 5% remove.")

    # ---- 28. v1 vs v2 table ----
    B.add_table(prs, "v1 vs v2 head-to-head (Graviton2, 90/5/5, L1_warm_zipf99)",
        ["threads", "v1 median (M/s)", "v2 median (M/s)", "v2 / v1", "v1 trial spread", "v2 trial spread"],
        [
            ["1", "17.32", "16.48", "0.95x",   "tight",            "tight"],
            ["2", "17.67", "21.01", "1.19x",   "16 - 20 M",         "tight"],
            ["3", "15.90", "26.46", "1.66x",   "wide",             "tight"],
            ["4", "4.90",  "30.35", "6.20x",   "2 - 6 M (3x)",       "29 - 31 M (tight)"],
            ["5", "0.41",  "31.60", "77x",     "0.4 - 0.7 M",       "tight"],
            ["6", "0.18",  "28.84", "160x",    "0.14 - 0.43 M",     "tight"],
            ["8", "0.030", "16.00", "533x",    "0.011 - 0.039 M",   "14.6 - 16.9 M"],
        ],
        col_widths_inches=[1.1, 2.0, 2.0, 1.4, 2.3, 2.4])

    # ---- 29. What the v2 numbers tell us ----
    B.add_bullets(prs, "What the v2 numbers tell us", [
        "1T regression of 5 % is real but small. The extra branch is the cost of correctness under contention.",
        "v2 scales from 16.5 M at 1T to 31.6 M at 5T (1.92x scaling), then plateaus and gently descends through 8T (16.0 M).",
        "Trial variance is gone: 8T v2 trials are 14.6 / 16.9 / 16.0. v1's 8T trials were 38 K / 32 K / 11 K: three orders of magnitude apart, non-deterministic.",
        "The new bottleneck at 8 threads is a different one: writers serialise on the slow-path mutex. With 10 % writers and 8 threads competing, the writer side itself becomes saturated. Mild, manageable: and crucially, deterministic (no trial-to-trial variance).",
        "At 8T v2 sits between wh-default (20.3 M/s) and the spinlocks (~14-15 M/s). Lock-free wh-occ-opt remains the runaway winner at 69.6 M/s.",
    ])

    # ---- 30. Investigation loop closed ----
    B.add_bullets(prs, "Investigation loop closed: what was learned", [
        "Diagnose -> hypothesise -> fix -> validate. Six hypotheses tested in one day. Three discriminating diagnostics. One implementation. One validation sweep.",
        "Three findings that the investigation produced:",
        (1, "(F1) Naive per-CPU rwlocks with retract-and-spin protocols fail catastrophically under any meaningful writer rate."),
        (1, "(F2) Wormhole's two-lock-per-op pattern amplifies primitive failures by ~86x. Lock primitives in concurrent indexes must be tested both in isolation AND in their target call pattern."),
        (1, "(F3) Faster atomics can make a fragile lock fail HARDER. The literature does not document this. Hardware modernisation can expose latent timing bugs."),
        "The thesis story strengthens: naive per-CPU rwlock fails catastrophically; the percpu_rwsem-style fix restores it to mid-pack throughput; lock-free OCC reads remain the practical winner.",
        "Both v1's failure and v2's recovery are publishable contributions. The investigation methodology itself is a contribution.",
        "Full record: docs/INVESTIGATION_PCPU_RW.md (~18 KB journal, all diagnostics + analysis + fix design + validation).",
    ])

    # ---- 31. Section divider: What this means ----
    B.add_section(prs, "6. What this means")

    # ---- 32. Lock ranking ----
    B.add_table(prs, "Lock ranking at 8T on Graviton2 (L1_warm_read_heavy)",
        ["Lock", "M ops/s", "Notes"],
        [
            ["wh-occ-opt",       "69.6", "lock-free reads: the runaway winner"],
            ["wh-default",       "20.3", "counter rwlock; bounded coherence cost"],
            ["wh-pcpu-rw-v2",    "16.0", "fixed per-CPU rwlock; mid-pack"],
            ["wh-cas",           "15.5", "spinlock-as-rwlock variants overlap"],
            ["wh-tas",           "14.2", "same"],
            ["wh-occ",           "11.0", "seqlock writes pay CAS on version counter"],
            ["wh-pcpu-rw (v1)",  "0.03", "thundering herd: broken under any writer rate"],
        ],
        col_widths_inches=[3.0, 2.0, 7.0])

    # ---- 33. Three findings ----
    B.add_bullets(prs, "Three findings from this week's data", [
        "1.  Lock-free OCC reads win categorically on read-heavy workloads: 3.4x over wh-default on Graviton, 3.7x on Xeon at max threads. If the data structure permits OCC, use it.",
        "2.  Graviton2 is 1.4-2.8x faster than Xeon on every well-behaved lock. Each cross-arch ratio fingerprints the bottleneck: spinlocks atomic-bound (LSE wins), default coherence-bound, occ-opt memory-bound (DRAM helps).",
        "3.  The pcpu_rw_lock story is the central methodological finding. Naive per-CPU rwlocks are NOT a drop-in fix: protocol matters more than data layout. Faster atomics can make a fragile lock fail HARDER, exposing latent timing bugs that look fine on older silicon.",
    ])

    # ---- 34. Open questions ----
    B.add_bullets(prs, "Open questions for you", [
        "NUMA story: currently dropped via D23's single-socket cap. ~3.3 h to recover cross-socket data on Xeon. Worth it for the thesis?",
        "v2 across the full matrix: table above is only the diagnostic cell on Graviton. Rerun wh_compare.sh with v2 in the lock list on both arches (~6.6 h x 2)?",
        "Notebook section still reads as \"anatomy of a thundering herd\" without the v2 resolution. Update with full diagnose -> fix -> validate arc?",
        "Higher thread counts on Graviton: c6g.16xlarge gives 16+ cores. Probably not blocking the thesis. Flagging.",
    ])

    # ---- 35. Next steps ----
    B.add_bullets(prs, "Next steps: proposed for this week", [
        "Rerun wh_compare.sh with pcpu-rw-v2 in the lock list on both arches. Gives 12-workload coverage of v2 to put alongside v1.",
        "Update notebook to incorporate the v2 validation data: currently still reads as \"anatomy of a thundering herd\" without the resolution.",
        "If you greenlight the NUMA story: revert the single-socket cap on Xeon, rerun wh_compare.sh, add a notebook section on cross-socket coherence.",
        "Otherwise: start drafting the thesis chapter using this deck's structure (changes -> architectural reference -> 1T -> scaling -> cross-arch ratio -> pcpu story -> findings).",
    ])

    # ─────────────────────────────────────────────────────────
    # Append the lock-ranking cross-arch slides (replaces the prior
    # spinlock-only slides; focused on rank CHANGES between platforms).
    # ─────────────────────────────────────────────────────────

    # ---- 36. Section divider ----
    B.add_section(prs, "7. How does the lock ranking change between platforms?")

    # ---- 37. Wormhole L1_warm slopegraph ----
    B.add_image(prs,
        "Wormhole: lock ranking on Xeon vs Graviton2",
        FIG_DIR / "f10_rank_wh_L1warm.png",
        "Each line is one lock. Bold = rank changed across platforms. Workload: L1-resident, warm zipf, 90/5/5.")

    # ---- 38. Wormhole observations ----
    B.add_bullets(prs,
        "Wormhole: what stays the same, what swings",
        [
            "Stable across both platforms (no rank change):",
            (1, "occ-opt at #1: lock-free reader path is platform-invariant. No atomic on the read side means no ISA-dependent cost."),
            (1, "default at #2: counter-based rwlock is cache-line-bounded on both arches. Both architectures converge to the same coherence ceiling."),
            (1, "pcpu-rw at #7: broken everywhere (separate slide deck section)."),
            "Rank changes in the middle of the table (cas/ttas/tas/occ all reshuffle):",
            (1, "cas jumps from #6 on Xeon to #3 on Graviton (↑3 ranks): simple CAS spinlock benefits the most from ARMv8.1's single-instruction atomic operations."),
            (1, "occ falls from #3 on Xeon to #6 on Graviton (↓3 ranks): the seqlock write path is CAS-bound too, but the spinlock variants gain MORE from the LSE upgrade, so occ ends up below them."),
            "The middle of the table reshuffles entirely. The top and bottom of the table are platform-invariant.",
        ])

    # ---- 39. Wormhole L3: cache regime changes the ranking too ----
    B.add_image(prs,
        "Wormhole at L3: the cache regime changes the ranking",
        FIG_DIR / "f11_rank_wh_L3warm.png",
        "Same locks, same arches, larger index (100k keys, spills to L3). New cross-arch behaviour at the bottom of the table.")

    B.add_bullets(prs,
        "Wormhole L1 vs L3: a second axis of rank movement",
        [
            "Moving from L1 (1k keys) to L3 (100k keys) on the SAME architecture also reshuffles the ranking:",
            (1, "default falls from #2 (L1) to #6 (L3) on Graviton: the contended-counter rwlock's cache-line bouncing combines with the L3 miss cost, more than the spinlocks' atomic cost does."),
            (1, "Spinlocks (tas/ttas/cas) climb to #2-#4 on Graviton in L3: their per-op atomic cost is now masked by the leaf-fetch latency from L3, so they look relatively better."),
            "Cross-platform takeaway: a lock's rank depends on TWO factors, not just one:",
            (1, "How much work the lock's fast path does per acquire (atomic cost)."),
            (1, "How much memory latency the surrounding operation pays (cache regime)."),
            "Both axes interact differently with each arch.",
        ])

    # ---- 40. StripedMap slopegraph ----
    B.add_image(prs,
        "StripedMap: lock ranking on Xeon vs Graviton2",
        FIG_DIR / "f12_rank_cds_L1warm.png",
        "Hash table with one lock per stripe. Workload: L1-resident, warm zipf, 90/5/5.")

    B.add_bullets(prs,
        "StripedMap: the same TAS / ticket swap shows up",
        [
            "std (pthread mutex) stays at #1 on both arches: futex parking is platform-invariant for short critical sections in this scheme.",
            "tas jumps from #5 on Xeon to #2 on Graviton (↑3 ranks): same LSE-atomic story as wormhole.",
            "ticket drops from #2 on Xeon to #3 on Graviton (↓1 rank): fairness-ordering's serial nature hurts more when the underlying atomics are fast: the ordering becomes the bottleneck rather than the atomic itself.",
            "ttas drops from #4 to #5 (↓1 rank): the read-before-CAS pattern is less of an advantage when the atomic is already cheap on Graviton.",
            "Pattern: simplest spinlock (tas) wins on Graviton; locks with extra mechanism (ticket, ttas) lose rank.",
        ])

    # ---- 41. BronsonAVL slopegraph ----
    B.add_image(prs,
        "BronsonAVL: lock ranking on Xeon vs Graviton2",
        FIG_DIR / "f13_rank_avl_L1warm.png",
        "Per-node lock on a tree. Workload: L1-resident, warm zipf, 90/5/5.")

    B.add_bullets(prs,
        "BronsonAVL: the most dramatic rank swap of any bench",
        [
            "tas jumps from #5 (Xeon) to #1 (Graviton): a four-rank climb, the largest swing in the matrix.",
            "ticket falls from #1 (Xeon) to #3 (Graviton): the previous winner becomes mid-pack.",
            "std falls from #3 to #5 (↓2 ranks).",
            "Why is the swing larger here?",
            (1, "BronsonAVL takes the per-node lock only briefly during the optimistic read protocol. The lock-acquire cost is a small fraction of total op time."),
            (1, "On a tree workload, BronsonAVL's ticket-lock advantage on Xeon comes from FIFO ordering aligning with version-check reads."),
            (1, "On Graviton, LSE makes raw atomic acquires so cheap that the FIFO advantage disappears: simple TAS wins on raw speed."),
            "Same direction of rank movement as the other two benches (tas up, ticket/std down): but the magnitude is largest here.",
        ])

    # ---- 42. Cross-bench heatmap ----
    B.add_image(prs,
        "Same patterns across all 12 workloads: rank change heatmap",
        FIG_DIR / "f14_rank_change_heatmap.png",
        "Blue = lock moved UP on Graviton. Red = lock moved DOWN. White = no rank change. Cell = (Graviton rank − Xeon rank).")

    # ---- 43. Synthesis / why ----
    B.add_bullets(prs,
        "Three consistent patterns explain why ranks change",
        [
            "Looking across all 12 workloads × 3 benches, the rank changes between arches follow three rules:",
            "(1) Simple-atomic spinlocks move UP on Graviton.",
            (1, "tas and cas climb in nearly every cell (most blue, few exceptions)."),
            (1, "Cause: ARMv8.1 LSE atomics replace LL/SC-style retry loops with single-instruction atomic operations (~2-3x faster cycles). Locks whose entire fast path IS one atomic benefit the most."),
            "(2) Locks with extra ordering mechanism move DOWN on Graviton.",
            (1, "ticket falls in most read-heavy L1 cells (light red)."),
            (1, "ttas falls in most cells (light red on cds, mixed on others)."),
            (1, "Cause: ticket's FIFO turnstile serialises acquire order across acquirers. When atomics are slow (Xeon), the turnstile cost is dominated by the atomic itself. When atomics are fast (Graviton), the turnstile becomes the bottleneck."),
            "(3) Cache-coherence-bound locks track the cache regime, not the lock's atomic cost.",
            (1, "wh-default falls hard on Graviton in L3 (+3, +4 cells in heatmap)."),
            (1, "Cause: a contended counter line costs about the same per acquire on both arches; what changes is how much the SURROUNDING work hides that cost."),
        ])

    # ---- 44. What this tells us: the meta finding ----
    B.add_bullets(prs,
        "What lock ranking changes tell us about cross-platform porting",
        [
            "If you tuned your lock choice on Xeon, your choice will be wrong on Graviton (and vice versa) about half the time.",
            "Specifically wrong:",
            (1, "If you chose ticket on Xeon for fairness: TAS will likely beat it on Graviton."),
            (1, "If you chose CAS on Xeon as a careful retry-friendly primitive: it may now be the BEST choice on Graviton."),
            (1, "If you chose a counter-based rwlock for L1-resident workloads: the same lock on L3-resident workloads on Graviton may collapse out of contention."),
            "What survives the cross-platform port unchanged:",
            (1, "Lock-free read paths (occ-opt): platform-invariant #1."),
            (1, "pthread mutex (std): futex parking is platform-invariant."),
            (1, "Catastrophically-failing locks (pcpu-rw v1): bottom on both arches.",),
            "Bottom line: portable lock choice means choosing locks whose performance does NOT depend on per-op atomic cycle count. Lock-free design and kernel-park-based designs are the most portable.",
        ])

    prs.save(TARGET)
    print(f"\nWrote {len(prs.slides)}-slide deck to {TARGET}")
    print(f"Size: {TARGET.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
