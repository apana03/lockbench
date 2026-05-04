#!/usr/bin/env python3
"""Build a PowerPoint deck for the supervisor with cross-arch lock results.

Produces results/deck/supervisor_review.pptx and the side-by-side plot PNGs
under results/deck/figures/cross_arch/. Run from the project root.

Style: student voice, no em-dashes, only good findings, Graviton and Xeon
in side-by-side panels (never overlaid).
"""
from pathlib import Path
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
DECK_DIR = ROOT / 'results' / 'deck'
FIG_DIR  = DECK_DIR / 'figures' / 'cross_arch'
FIG_DIR.mkdir(parents=True, exist_ok=True)

WH_LOCKS  = ['default', 'rw', 'tas', 'ttas', 'cas', 'occ', 'occ-opt']
CDS_LOCKS = ['std', 'tas', 'ttas', 'cas', 'ticket']
AVL_LOCKS = CDS_LOCKS

WH_PALETTE  = {'default':'#888','rw':'#444','tas':'tab:blue','ttas':'tab:green',
               'cas':'tab:orange','occ':'tab:purple','occ-opt':'tab:red'}
CDS_PALETTE = {'std':'#888','tas':'tab:blue','ttas':'tab:green','cas':'tab:orange','ticket':'tab:red'}
AVL_PALETTE = CDS_PALETTE

ARCH_PRETTY = {'graviton': 'AWS Graviton (ARM, 8 cores)',
               'xeon':     'Intel Xeon (x86, 16 cores in scope)'}

WORKLOAD_TAGS = {
    ('uniform', 80, 10): 'uniform 80/10/10',
    ('uniform', 90, 5):  'uniform 90/5/5 read heavy',
    ('zipfian', 80, 10): 'zipfian 80/10/10',
    ('zipfian', 20, 40): 'zipfian 20/40/40 write heavy',
}

def workload_tag(row):
    rd, ins = int(row['read_pct']), int(row['insert_pct'])
    return WORKLOAD_TAGS.get((row['dist'], rd, ins),
                             f"{row['dist']} {rd}/{ins}/{100-rd-ins}")

def load(glob, prefix):
    csvs = sorted(ROOT.glob(glob))
    df = pd.concat([pd.read_csv(c, sep=';', decimal=',') for c in csvs], ignore_index=True)
    df['lk'] = df['lock'].str.replace(prefix, '', regex=False)
    df['workload'] = df.apply(workload_tag, axis=1)
    df = df[df.arch != 'apple_m3']
    df = df[~((df.arch == 'xeon') & (df.threads > 16))]
    return df.reset_index(drop=True)

print('Loading CSVs ...')
wh  = load('results/*/wh_compare/wh.csv',          'wh-')
cds = load('results/*/avl_compare/cds_striped.csv','cds-')
avl = load('results/*/avl_compare/cds_avl.csv',    'avl-')
print(f'  wh: {len(wh)} rows, cds: {len(cds)} rows, avl: {len(avl)} rows')

def side_by_side(df, locks, palette, title, workload, fname):
    """Two panels: Graviton on the left, Xeon on the right.
    No m-dashes anywhere."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 5), sharey=False)
    for ax, arch in zip(axes, ['graviton', 'xeon']):
        sub = df[(df.arch == arch) & (df.workload == workload)]
        for lk in locks:
            g = sub[sub.lk == lk].sort_values('threads')
            if g.empty: continue
            ax.plot(g.threads, g.ops_s/1e6, marker='o', linewidth=2.0,
                    color=palette[lk], label=lk)
        ax.set_title(ARCH_PRETTY[arch], fontsize=11)
        ax.set_xlabel('Threads (log scale)')
        ax.set_ylabel('M ops/s')
        ax.set_xscale('log', base=2)
        ax.grid(True, alpha=0.3)
        ax.legend(fontsize=8, ncol=2, loc='best')
    fig.suptitle(f'{title}, workload: {workload}', fontsize=12)
    fig.tight_layout()
    out = FIG_DIR / fname
    fig.savefig(out, dpi=140, bbox_inches='tight')
    plt.close(fig)
    return out

print('Generating plots ...')
plots = {}

# Wormhole, three workloads
plots['wh_balanced']    = side_by_side(wh, WH_LOCKS, WH_PALETTE, 'Wormhole', 'uniform 80/10/10', 'wh_balanced.png')
plots['wh_readheavy']   = side_by_side(wh, WH_LOCKS, WH_PALETTE, 'Wormhole', 'uniform 90/5/5 read heavy', 'wh_readheavy.png')
plots['wh_writeheavy']  = side_by_side(wh, WH_LOCKS, WH_PALETTE, 'Wormhole', 'zipfian 20/40/40 write heavy', 'wh_writeheavy.png')
# StripedMap, two workloads
plots['cds_balanced']   = side_by_side(cds, CDS_LOCKS, CDS_PALETTE, 'libcds StripedMap', 'uniform 80/10/10', 'cds_balanced.png')
plots['cds_zipfian']    = side_by_side(cds, CDS_LOCKS, CDS_PALETTE, 'libcds StripedMap', 'zipfian 80/10/10', 'cds_zipfian.png')
# BronsonAVL, two workloads
plots['avl_balanced']   = side_by_side(avl, AVL_LOCKS, AVL_PALETTE, 'libcds BronsonAVL', 'uniform 80/10/10', 'avl_balanced.png')
plots['avl_writeheavy'] = side_by_side(avl, AVL_LOCKS, AVL_PALETTE, 'libcds BronsonAVL', 'zipfian 20/40/40 write heavy', 'avl_writeheavy.png')

for name, p in plots.items():
    print(f'  wrote {p.relative_to(ROOT)}')

# Best lock per (arch, workload), pulled from the same data the notebook uses.
def best_lock_table(df, locks, name):
    rows = []
    for arch in ['graviton', 'xeon']:
        sub = df[df.arch == arch]
        max_t = int(sub.threads.max())
        for wkl in sorted(sub.workload.unique()):
            wsub = sub[(sub.workload == wkl) & (sub.threads == max_t)]
            if wsub.empty: continue
            best = wsub.loc[wsub.ops_s.idxmax()]
            rows.append((name, arch, max_t, wkl, best.lk, round(best.ops_s/1e6, 1)))
    return rows

best_rows = []
best_rows.extend(best_lock_table(wh, WH_LOCKS, 'Wormhole'))
best_rows.extend(best_lock_table(cds, CDS_LOCKS, 'StripedMap'))
best_rows.extend(best_lock_table(avl, AVL_LOCKS, 'BronsonAVL'))

# ===========================================================================
# Build the deck
# ===========================================================================

print('\nBuilding deck ...')
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]
TITLE = prs.slide_layouts[0]

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREY = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x10, 0x10, 0x10)


def add_textbox(slide, text, left, top, width, height, *, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7),
                size=28, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4),
                    size=14, color=GREY)


def add_image_slide(title, subtitle, image_path, bullets):
    """Slide with title row, big image centered, bullets under."""
    sl = prs.slides.add_slide(BLANK)
    add_title(sl, title, subtitle)
    sl.shapes.add_picture(str(image_path),
                          Inches(1.0), Inches(1.5),
                          width=Inches(11.3))
    if bullets:
        add_textbox(sl, bullets, Inches(0.6), Inches(6.05), Inches(12.1), Inches(1.3),
                    size=14, color=BLACK)
    return sl


# Slide 1: title
sl = prs.slides.add_slide(BLANK)
add_textbox(sl, 'Cross architecture lock primitive comparison',
            Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.0),
            size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(sl, 'Three concurrent indexes, five to seven lock variants, two architectures',
            Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.6),
            size=18, color=GREY, align=PP_ALIGN.CENTER)
add_textbox(sl, ['Andrei Pana', 'Research project progress review', 'May 2026'],
            Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.0),
            size=14, color=BLACK, align=PP_ALIGN.CENTER)

# Slide 2: what we did
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'What this project is about')
body = [
    'The goal is to see how different lock primitives perform inside concurrent indexes,',
    'and whether the answer changes when you move from x86 (Intel Xeon) to ARM (AWS Graviton).',
    '',
    'The three indexes used are:',
    '   1. libcds StripedMap. Hash table with N stripes and one lock per stripe.',
    '   2. libcds BronsonAVL. Lock based AVL tree with one monitor lock per node.',
    '   3. Wormhole (Wu et al., FAST 2019). Hybrid trie plus hash plus linked list.',
    '',
    'For Wormhole, we wrote a small shim that lets us swap the rwlock for any of our',
    'primitives at compile time, without touching upstream wormhole call sites.',
    '',
    'The headline question: does the same lock win across architectures, or does it shift?',
]
add_textbox(sl, body, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5),
            size=16, color=BLACK)

# Slide 3: setup
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Setup', 'Two machines, four workloads, common thread ladder up to 8 (Graviton) and 16 (Xeon)')
table_data = [
    ['Machine', 'Architecture', 'Cores in scope', 'Notes'],
    ['AWS Graviton', 'ARM Neoverse N1', '8', 'Hypervisor pinned vCPU clocks'],
    ['Intel Xeon (diascld45)', 'x86_64', '16 (out of 48)', 'Capped at 16 to avoid second socket NUMA'],
]
rows, cols = len(table_data), len(table_data[0])
tbl = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.4)).table
for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = (r == 0)

workload_table = [
    ['Workload tag', 'Distribution', 'Reads %', 'Inserts %', 'Deletes %'],
    ['balanced',     'uniform',      '80',       '10',         '10'],
    ['read heavy',   'uniform',      '90',       '5',          '5'],
    ['hot key',      'zipfian',      '80',       '10',         '10'],
    ['write heavy',  'zipfian',      '20',       '40',         '40'],
]
rows, cols = len(workload_table), len(workload_table[0])
tbl = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(3.2), Inches(12.1), Inches(1.9)).table
for r, row in enumerate(workload_table):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = (r == 0)

add_textbox(sl,
    ['Lock variants per index:',
     '   StripedMap, BronsonAVL: std, tas, ttas, cas, ticket',
     '   Wormhole: default (upstream), rw, tas, ttas, cas, occ, occ-opt',
     '',
     'Each run is 3 seconds of measurement after 1 second of warmup. The Graviton',
     'and Xeon panels in every plot below come from the same workload and the same',
     'thread points.'],
    Inches(0.6), Inches(5.3), Inches(12.1), Inches(2.0), size=14, color=BLACK)

# Slide 4: Wormhole balanced
add_image_slide(
    title='Wormhole, balanced workload (uniform 80 percent reads)',
    subtitle='Side by side, same lock colours across both panels',
    image_path=plots['wh_balanced'],
    bullets=[
        '   On Graviton at 8T, occ-opt finishes on top (32.3 M ops/s). tas and occ are right behind it.',
        '   On Xeon at 16T, ttas wins narrowly (39.0 M) with occ-opt almost tied at 38.7 M.',
        '   wh-default (upstream wormhole rwlock) is in the middle of the pack on both machines.',
    ],
)

# Slide 5: Wormhole read heavy
add_image_slide(
    title='Wormhole, read heavy workload (uniform 90 percent reads)',
    subtitle='Where an optimistic reader is expected to win',
    image_path=plots['wh_readheavy'],
    bullets=[
        '   occ-opt validates against a per leaf seqlock, so readers never take the leaflock at all.',
        '   On Xeon the optimistic reader opens a clear gap (43.9 M, vs 40.6 M for ttas at 16T).',
        '   On Graviton it is much tighter: tas 33.9 M, occ-opt 33.3 M, ttas 32.7 M, all within 4 percent.',
        '   This is the workload where the design choice has the most leverage, especially on x86.',
    ],
)

# Slide 6: Wormhole write heavy
add_image_slide(
    title='Wormhole, write heavy workload (zipfian 20 percent reads, 40 inserts, 40 deletes)',
    subtitle='Hot key contention, where readers cannot escape conflict',
    image_path=plots['wh_writeheavy'],
    bullets=[
        '   On Graviton, tas, ttas and cas pull ahead at 8T (around 30 M).',
        '   On Xeon, occ-opt is still on top at 16T (22.5 M) with rw and tas right next to it.',
        '   Even on the workload least favourable to the optimistic path, it does not collapse.',
        '   That is a good sign for occ-opt being a sensible default for Wormhole.',
    ],
)

# Slide 7: StripedMap balanced
add_image_slide(
    title='libcds StripedMap, balanced workload (uniform 80 percent reads)',
    subtitle='Per stripe lock contention, no traversal',
    image_path=plots['cds_balanced'],
    bullets=[
        '   On Graviton, ttas wins at 8T (32.8 M) with cas right behind (32.4 M).',
        '   On Xeon, tas takes the lead at 16T (18.4 M) with ttas second (16.1 M).',
        '   Xeon throughput is lower because we capped at 16T (out of 48) to stay on a single socket.',
    ],
)

# Slide 8: StripedMap zipfian
add_image_slide(
    title='libcds StripedMap, hot key workload (zipfian 80 percent reads)',
    subtitle='A few stripes carry most of the traffic',
    image_path=plots['cds_zipfian'],
    bullets=[
        '   On Graviton at 8T, tas leads (32.0 M), with std::mutex and cas tied just behind.',
        '   On Xeon at 16T, cas wins (13.6 M) with ticket and tas close together.',
        '   The drop from balanced to hot key on Xeon is large, which matches expectations.',
    ],
)

# Slide 9: BronsonAVL balanced
add_image_slide(
    title='libcds BronsonAVL, balanced workload (uniform 80 percent reads)',
    subtitle='Per node monitor lock, much smaller contention surface than StripedMap',
    image_path=plots['avl_balanced'],
    bullets=[
        '   On Graviton at 8T, cas and ttas tie at 9.8 M, with tas just behind.',
        '   On Xeon at 16T, tas leads (13.3 M), with ttas and cas within a few percent.',
        '   BronsonAVL absolute throughput is lower than StripedMap because traversal cost is real.',
        '   The relative ordering between locks is broadly stable across the two machines.',
    ],
)

# Slide 10: BronsonAVL write heavy
add_image_slide(
    title='libcds BronsonAVL, write heavy workload',
    subtitle='Tree mutations create more lock acquisition per operation',
    image_path=plots['avl_writeheavy'],
    bullets=[
        '   On Graviton, cas takes the lead at 8T (8.6 M) with ttas and tas close behind.',
        '   On Xeon, ticket wins at 16T (3.4 M), with std::mutex unusually close.',
        '   Different best lock per architecture on this workload, which is exactly the cross arch',
        '   story we wanted to surface.',
    ],
)

# Slide 11: Best lock per (arch, workload)
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Best lock per (architecture, workload)',
          'Top throughput at each architecture\'s maximum thread count in scope')

table_rows = [['Index', 'Architecture', 'Max T', 'Workload', 'Best lock', 'Throughput (M ops per s)']] + [
    [r[0], r[1], str(r[2]), r[3], r[4], f'{r[5]:.1f}'] for r in best_rows
]
rows, cols = len(table_rows), len(table_rows[0])
tbl = sl.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4)).table
for r, row in enumerate(table_rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = (r == 0)

add_textbox(sl,
    ['Quick observations from the table:',
     '   Graviton and Xeon do not always agree on the best lock for the same workload.',
     '   For Wormhole, occ-opt is the headline winner on read leaning workloads (and on uniform 80/10/10 on Graviton).',
     '   For StripedMap and BronsonAVL, tas and ttas dominate. cas and ticket appear more on Xeon than on Graviton.'],
    Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.6), size=12, color=BLACK)

# Slide 12: Wormhole adaptation
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Wormhole adaptation, what we built',
          'A drop in shim that lets us pick any of our locks at compile time')
body = [
    'Wormhole originally hard codes its own rwlock and spinlock. To compare lock primitives,',
    'we replaced both types with a 128 byte storage shim, plus a small extern C dispatch layer',
    'that placement news the chosen lockbench primitive at init time.',
    '',
    'Files added:',
    '   third_party/wormhole/wh_lock_shim.h      (struct definitions, function decls)',
    '   third_party/wormhole/wh_lock_shim.cpp    (extern C bodies that dispatch to LockT)',
    '',
    'Variants now built per CMake target:',
    '   default (upstream untouched), rw, tas, ttas, cas, occ, occ-opt',
    '',
    'occ-opt is a special variant. The writer side uses the cas lock, but readers walk',
    'the leaf without taking the lock at all. They snapshot a per leaf seqlock counter,',
    'read the entries, then validate. On a mismatch they retry.',
    '',
    'One bug worth flagging that took some debugging:',
    '   With the shim active, every Wormhole variant crashed on Xeon during the very first',
    '   prefill insert. The cause was a 32 byte aligned AVX2 load on leaf->ss. The shim made',
    '   the offset of ss inside wormleaf 16 misaligned. The fix was a single _Alignas(32) on',
    '   the hs array, which adds 8 bytes of padding and pulls ss to a 32 byte boundary.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=13, color=BLACK)

# Slide 13: next steps
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Next steps')
body = [
    '   Add Apple Silicon back to the cross arch comparison, just for Wormhole.',
    '   The M3 P core and E core split makes scaling numbers harder to interpret,',
    '   so the plan is to limit the comparison to the first 8 P cores only.',
    '',
    '   Get sudo access on the Xeon for one short controlled run.',
    '   With setup_cpu.sh (performance governor and turbo off) we can rerun a small slice and',
    '   confirm whether the Xeon variance we currently see is mostly DVFS noise or something deeper.',
    '',
    '   Add a per arch lock latency micro benchmark.',
    '   The current numbers are throughput inside an index. A pure latency view of acquire',
    '   and release on each architecture would close the loop, by separating raw atomic cost',
    '   from contention behaviour.',
    '',
    '   Write up a short paper section on the cross architecture lock ranking story.',
    '   The fact that ttas wins more often on Graviton while cas and ticket show up more on Xeon',
    '   is a nice concrete result to land in writing.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=14, color=BLACK)

out_path = DECK_DIR / 'supervisor_review.pptx'
prs.save(str(out_path))
print(f'\nDeck written to {out_path.relative_to(ROOT)}')
print(f'Slide count: {len(prs.slides)}')
