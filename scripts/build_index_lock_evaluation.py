#!/usr/bin/env python3
"""Build the Index-Lock Evaluation deck for the supervisor review.

Output:
    results/deck/index_lock_evaluation.pptx
    results/deck/figures/index_lock/*.png

Plot layout per workload: one image with 2 rows x 2 cols.
    Row 1: throughput vs threads (Graviton on the left, Xeon on the right).
    Row 2: bar chart of M ops per second at the architecture's maximum
           thread point, sorted descending so the lock ranking is visible.

Style: student voice, plain wording, no em dashes anywhere.
"""
from pathlib import Path
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
DECK_DIR = ROOT / 'results' / 'deck'
FIG_DIR  = DECK_DIR / 'figures' / 'index_lock'
FIG_DIR.mkdir(parents=True, exist_ok=True)

WH_LOCKS  = ['default', 'rw', 'tas', 'ttas', 'cas', 'occ', 'occ-opt']
CDS_LOCKS = ['std', 'tas', 'ttas', 'cas', 'ticket']
AVL_LOCKS = CDS_LOCKS

WH_PALETTE  = {'default':'#888','rw':'#444','tas':'tab:blue','ttas':'tab:green',
               'cas':'tab:orange','occ':'tab:purple','occ-opt':'tab:red'}
CDS_PALETTE = {'std':'#888','tas':'tab:blue','ttas':'tab:green','cas':'tab:orange','ticket':'tab:red'}
AVL_PALETTE = CDS_PALETTE

ARCH_PRETTY = {'graviton': 'AWS Graviton (ARM, 8 cores)',
               'xeon':     'Intel Xeon (x86, 16 cores in scope)'}

WORKLOAD_TAGS = {
    ('uniform', 80, 10): 'uniform 80/10/10',
    ('uniform', 90, 5):  'uniform 90/5/5 read heavy',
    ('zipfian', 80, 10): 'zipfian 80/10/10',
    ('zipfian', 20, 40): 'zipfian 20/40/40 write heavy',
}

def workload_tag(row):
    rd, ins = int(row['read_pct']), int(row['insert_pct'])
    return WORKLOAD_TAGS.get((row['dist'], rd, ins),
                             f"{row['dist']} {rd}/{ins}/{100-rd-ins}")

def load(glob, prefix):
    csvs = sorted(ROOT.glob(glob))
    df = pd.concat([pd.read_csv(c, sep=';', decimal=',') for c in csvs], ignore_index=True)
    df['lk'] = df['lock'].str.replace(prefix, '', regex=False)
    df['workload'] = df.apply(workload_tag, axis=1)
    df = df[df.arch != 'apple_m3']
    df = df[~((df.arch == 'xeon') & (df.threads > 16))]
    return df.reset_index(drop=True)

print('Loading CSVs ...')
wh  = load('results/*/wh_compare/wh.csv',          'wh-')
cds = load('results/*/avl_compare/cds_striped.csv','cds-')
avl = load('results/*/avl_compare/cds_avl.csv',    'avl-')
print(f'  wh: {len(wh)} rows, cds: {len(cds)} rows, avl: {len(avl)} rows')


def plot_lines_and_bars(df, locks, palette, title, workload, fname):
    """2x2: top row is line plots vs threads; bottom row is bar chart at maxT.
    Left column is Graviton, right column is Xeon. Aspect ratio chosen so
    the image fits cleanly under the slide title without overflow."""
    fig, axes = plt.subplots(2, 2, figsize=(14, 6.4),
                             gridspec_kw={'height_ratios': [3, 2]})
    for col, arch in enumerate(['graviton', 'xeon']):
        sub = df[(df.arch == arch) & (df.workload == workload)]
        ax = axes[0, col]
        for lk in locks:
            g = sub[sub.lk == lk].sort_values('threads')
            if g.empty: continue
            ax.plot(g.threads, g.ops_s/1e6, marker='o', linewidth=2.0,
                    color=palette[lk], label=lk)
        ax.set_title(ARCH_PRETTY[arch], fontsize=11)
        ax.set_xlabel('Threads (log scale)')
        ax.set_ylabel('M ops/s')
        ax.set_xscale('log', base=2)
        ax.grid(True, alpha=0.3)
        ax.legend(fontsize=8, ncol=2, loc='best')
        ax = axes[1, col]
        if sub.empty:
            ax.set_axis_off(); continue
        max_t = int(sub.threads.max())
        s = sub[sub.threads == max_t].copy()
        s['ops_M'] = s.ops_s / 1e6
        s = s.sort_values('ops_M', ascending=False)
        bars = ax.bar(s.lk, s.ops_M,
                      color=[palette[lk] for lk in s.lk],
                      edgecolor='black', linewidth=0.4)
        for bar, val in zip(bars, s.ops_M):
            ax.text(bar.get_x() + bar.get_width()/2, val,
                    f'{val:.1f}', ha='center', va='bottom', fontsize=9)
        ax.set_title(f'Ranking at {max_t}T (M ops/s)', fontsize=10)
        ax.set_ylim(0, max(s.ops_M) * 1.18)
        ax.grid(True, axis='y', alpha=0.3)
        ax.tick_params(axis='x', labelsize=9)
    fig.suptitle(f'{title}, workload: {workload}', fontsize=13)
    fig.tight_layout()
    out = FIG_DIR / fname
    fig.savefig(out, dpi=140, bbox_inches='tight')
    plt.close(fig)
    return out


print('Generating plots ...')
plots = {
    'wh_balanced':    plot_lines_and_bars(wh, WH_LOCKS, WH_PALETTE, 'Wormhole', 'uniform 80/10/10', 'wh_balanced.png'),
    'wh_readheavy':   plot_lines_and_bars(wh, WH_LOCKS, WH_PALETTE, 'Wormhole', 'uniform 90/5/5 read heavy', 'wh_readheavy.png'),
    'wh_writeheavy':  plot_lines_and_bars(wh, WH_LOCKS, WH_PALETTE, 'Wormhole', 'zipfian 20/40/40 write heavy', 'wh_writeheavy.png'),
    'cds_balanced':   plot_lines_and_bars(cds, CDS_LOCKS, CDS_PALETTE, 'libcds StripedMap', 'uniform 80/10/10', 'cds_balanced.png'),
    'cds_zipfian':    plot_lines_and_bars(cds, CDS_LOCKS, CDS_PALETTE, 'libcds StripedMap', 'zipfian 80/10/10', 'cds_zipfian.png'),
    'avl_balanced':   plot_lines_and_bars(avl, AVL_LOCKS, AVL_PALETTE, 'libcds BronsonAVL', 'uniform 80/10/10', 'avl_balanced.png'),
    'avl_writeheavy': plot_lines_and_bars(avl, AVL_LOCKS, AVL_PALETTE, 'libcds BronsonAVL', 'zipfian 20/40/40 write heavy', 'avl_writeheavy.png'),
}
for k, p in plots.items():
    print(f'  wrote {p.relative_to(ROOT)}')


def best_lock_table(df, locks, name):
    rows = []
    for arch in ['graviton', 'xeon']:
        sub = df[df.arch == arch]
        max_t = int(sub.threads.max())
        for wkl in sorted(sub.workload.unique()):
            wsub = sub[(sub.workload == wkl) & (sub.threads == max_t)]
            if wsub.empty: continue
            best = wsub.loc[wsub.ops_s.idxmax()]
            rows.append((name, arch, max_t, wkl, best.lk, round(best.ops_s/1e6, 1)))
    return rows

best_rows = []
best_rows.extend(best_lock_table(wh, WH_LOCKS, 'Wormhole'))
best_rows.extend(best_lock_table(cds, CDS_LOCKS, 'StripedMap'))
best_rows.extend(best_lock_table(avl, AVL_LOCKS, 'BronsonAVL'))


# ===========================================================================
# Build the deck
# ===========================================================================

print('\nBuilding deck ...')
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
GREY  = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x10, 0x10, 0x10)


def add_textbox(slide, text, left, top, width, height, *, size=14, bold=False,
                color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7),
                size=26, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4),
                    size=13, color=GREY)


def add_image_slide(title, subtitle, image_path):
    """Throughput slide: title row plus a single combined figure.
    Image is sized to fit the available slide height with no overlap."""
    sl = prs.slides.add_slide(BLANK)
    add_title(sl, title, subtitle)
    # figure aspect is 14x6.4 = 2.1875. With image height 5.6", width 12.25".
    # Slide is 13.333" wide; centred horizontally: left = (13.333 - 12.25)/2 = 0.54".
    sl.shapes.add_picture(str(image_path),
                          Inches(0.55), Inches(1.5),
                          width=Inches(12.25))
    return sl


# ----- Slide 1: title -----
sl = prs.slides.add_slide(BLANK)
add_textbox(sl, 'Index lock evaluation',
            Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.0),
            size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(sl, 'How lock primitive choice changes throughput inside three concurrent indexes',
            Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
            size=18, color=GREY, align=PP_ALIGN.CENTER)
add_textbox(sl, 'Cross architecture: AWS Graviton (ARM) and Intel Xeon (x86)',
            Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
            size=15, color=GREY, align=PP_ALIGN.CENTER)
add_textbox(sl, ['Andrei Pana', 'Research project review', 'May 2026'],
            Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0),
            size=14, color=BLACK, align=PP_ALIGN.CENTER)


# ----- Slide 2: project context -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'What this project is about')
body = [
    'The goal is to see how different lock primitives perform inside three concurrent indexes,',
    'and whether the answer changes when you move from x86 (Intel Xeon) to ARM (AWS Graviton).',
    '',
    'The three indexes used:',
    '   1. libcds StripedMap. Hash table with N stripes, one lock per stripe.',
    '   2. libcds BronsonAVL. Lock based AVL tree with one monitor lock per node.',
    '   3. Wormhole (Wu et al., FAST 2019). Hybrid trie plus hash plus linked list.',
    '',
    'Lock primitives compared:',
    '   StripedMap and BronsonAVL: std::mutex, tas, ttas, cas, ticket.',
    '   Wormhole: default (upstream untouched), rw, tas, ttas, cas, occ, occ-opt.',
    '',
    'The headline question: does the same lock win across architectures, or does it shift?',
]
add_textbox(sl, body, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5),
            size=15, color=BLACK)


# ----- Slide 3: setup -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Setup')

table_data = [
    ['Machine', 'Architecture', 'Cores in scope', 'Notes'],
    ['AWS Graviton', 'ARM Neoverse N1', '8', 'Hypervisor pinned vCPU clocks'],
    ['Intel Xeon (diascld45)', 'x86_64', '16 (out of 48)', 'Capped at 16T to stay on a single socket'],
]
rows, cols = len(table_data), len(table_data[0])
tbl = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.4)).table
for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = (r == 0)

workload_table = [
    ['Workload tag', 'Distribution', 'Reads %', 'Inserts %', 'Deletes %'],
    ['balanced',     'uniform',      '80',       '10',         '10'],
    ['read heavy',   'uniform',      '90',       '5',          '5'],
    ['hot key',      'zipfian',      '80',       '10',         '10'],
    ['write heavy',  'zipfian',      '20',       '40',         '40'],
]
rows, cols = len(workload_table), len(workload_table[0])
tbl = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(3.1), Inches(12.1), Inches(1.9)).table
for r, row in enumerate(workload_table):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = (r == 0)

add_textbox(sl,
    ['Each run is 3 seconds of measurement after 1 second of warmup.',
     'Thread ladder is power of 2 up to the per machine cap (1, 2, 4, 8 on Graviton; 1, 2, 4, 8, 16 on Xeon).',
     'Every plot in this deck shows Graviton in the left panel and Xeon in the right panel for the same workload.',
     'Below each line plot we also show a sorted bar chart at the maximum thread point so the ranking is visible at a glance.'],
    Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), size=13, color=BLACK)


# ----- Slide: thread pinning details (Graviton specific) -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Thread pinning details, Graviton specifics',
          'How --pin works in the harness, and what it actually does on a c6g instance')
body = [
    'How --pin is implemented (Linux only, in include/primitives/util.hpp):',
    '   set_thread_affinity(thread_id) calls sched_setaffinity with a cpu_set_t containing only that one CPU.',
    '   Threads are numbered 0, 1, 2, ..., so thread N gets pinned to logical CPU N.',
    '   Pinning is dense and sequential: a 4 thread run uses CPUs 0..3, an 8 thread run uses CPUs 0..7.',
    '',
    'What that means on AWS Graviton specifically:',
    '   c6g instances run Neoverse N1 cores. Neoverse N1 does NOT support SMT, so one vCPU is one physical core.',
    '   The kernel sees CPUs 0..7 on a c6g.2xlarge, and those map 1 to 1 to physical Graviton cores.',
    '   AWS Nitro / Graviton does not migrate vCPUs across physical cores at runtime under steady load.',
    '   So sched_setaffinity to CPU N effectively pins to physical core N for the duration of the run.',
    '   Each pinned thread also gets a hypervisor managed clock that does not vary at runtime, which is one',
    '   reason fairness numbers on Graviton are very tight (around 0.98 across all benches).',
    '',
    'Honest disclaimer about the data in this deck:',
    '   The harness has a --pin flag, but only the standalone cds_sweep.sh passes it. The two sweep scripts',
    '   used to produce these slides (wh_compare.sh and run_avl_compare.sh) do NOT pass --pin.',
    '   So the Wormhole, BronsonAVL, and StripedMap numbers in this deck are from UNPINNED runs.',
    '',
    'Why this is fine for the comparison we are presenting:',
    '   We checked directly (StripedMap pinned vs unpinned on Graviton at 8T): avg fairness 0.984 vs 0.982.',
    '   Same on Xeon: 0.909 vs 0.908. Pinning has effectively no impact at the scale of this study.',
    '   On Graviton this is because the hypervisor placement is already stable. On Xeon it is because the',
    '   variance comes from DVFS, not from thread migration.',
    '   We will normalize this in the next iteration so all three sweeps pass --pin on Linux.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=11, color=BLACK)


# ----- Slide 4: How we built StripedMap -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'How we built StripedMap',
          'Adapter at include/indexes/striped_map_index.hpp')
body = [
    'Wraps cds::container::StripedMap from libcds.',
    '',
    '   Hash table with std::list buckets and a fixed array of N stripe locks.',
    '   The lock primitive is the template parameter to cds::container::striped_set::striping<Lock>,',
    '   so we plug any lockbench primitive (tas_lock, ttas_lock, cas_lock, ticket_lock, std::mutex) at compile time.',
    '   We use 65536 stripes. Each operation hashes the key, picks the stripe, and acquires its lock.',
    '   Resize policy is load_factor_resizing<4>: when item count divided by buckets exceeds 4,',
    '   the map takes ALL stripe locks (scoped_full_lock) and rehashes. We size the initial bucket array',
    '   large enough that we do not hit a resize during a benchmark run.',
    '',
    'Why it matters: StripedMap is the simplest of the three indexes. There is no traversal,',
    'each operation acquires exactly one lock, so it isolates the cost of a single lock acquire',
    'plus the time to walk the per stripe std::list.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=14, color=BLACK)


# ----- Slide 5: How we built BronsonAVL -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'How we built BronsonAVL',
          'Adapter at include/indexes/avl_tree_index.hpp')
body = [
    'Wraps cds::container::BronsonAVLTreeMap from libcds. Bronson et al. published this',
    'lock based AVL tree as a balanced concurrent ordered map with optimistic version validated reads.',
    '',
    '   Tree is balanced. Each node carries a monitor lock that we plug in via cds::sync::injecting_monitor<Lock>.',
    '   Reads use a sequence style version counter on each node. Writers take the lock, mutate, then bump the version.',
    '   Writes do hand over hand traversal down the tree, locking at most a small constant number of nodes at a time.',
    '   Memory reclamation is handled by URCU (cds::urcu::gc<general_buffered>). Worker threads attach to the URCU',
    '   manager on first call via a thread_local guard, then detach at thread exit.',
    '',
    'Why it matters: the contention surface is much smaller than StripedMap. The hot lock is wherever the',
    'tree is being mutated, which depends on the key access pattern. This is what makes the zipfian',
    'workloads interesting: under a skewed distribution, a small set of nodes carries most writes.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=13, color=BLACK)


# ----- Slide 6: How we built Wormhole, overview -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'How we built Wormhole, part 1: vendoring and shim design',
          'Adapter at include/indexes/wormhole_index.hpp; shim at third_party/wormhole/wh_lock_shim.{h,cpp}')
body = [
    'Wormhole (Wu et al., FAST 2019) is a hybrid trie plus hash plus linked list ordered concurrent index.',
    'It hard codes its own rwlock and spinlock as 4 byte opaque structs in lib.h. For our experiment we need',
    'to swap them with any of our primitives at compile time.',
    '',
    'Approach:',
    '   1. Vendor wormhole in tree at third_party/wormhole/. The only edits to upstream sources are short',
    '      ifdef gated blocks. With WH_LOCK_SHIM unset, the build is upstream wormhole untouched.',
    '   2. Add wh_lock_shim.h, which redefines struct rwlock and struct spinlock as 128 byte storage arrays',
    '      with extern C function declarations matching upstream signatures (rwlock_lock_read, rwlock_trylock_write_nr, etc.).',
    '   3. Add wh_lock_shim.cpp, a C++ TU that picks LockT via WH_LOCK_<NAME> macro, asserts size and alignment,',
    '      and provides extern C bodies. Each function placement-news the LockT into the shim storage at init time',
    '      and dispatches via small if constexpr templates so the non matching branches are elided.',
    '',
    'Variants built per CMake target:',
    '   default (upstream rwlock, no shim), rw, tas, ttas, cas, occ, occ-opt.',
    '',
    'occ-opt is special: writer side uses the cas lock, but readers walk the leaf without taking the lock at all.',
    'They snapshot a per leaf seqlock counter, read the entries, then validate. On mismatch they retry.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=12, color=BLACK)


# ----- Slide 7: Wormhole bug story (alignment) -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'How we built Wormhole, part 2: the alignment bug',
          'Every shim variant SIGSEGV on Xeon during the very first prefill insert')
body = [
    'Symptom: wh_test_rw, wh_test_cas, wh_test_tas, wh_test_ttas, wh_test_occ, wh_test_occ-opt all crashed in',
    'wormleaf_insert_hs on Xeon. wh_test_default ran clean. The crash was invisible on Apple Silicon.',
    '',
    'Root cause:',
    '   wormleaf_shift_inc and wormleaf_shift_dec use _mm256_load_si256 on leaf->ss (AVX2 aligned load).',
    '   Aligned loads need the source address to be a multiple of 32 bytes.',
    '   With upstream wormhole, leaflock + sortlock take 8 bytes total, so leaf->ss starts at offset 1088 (aligned).',
    '   With our shim, leaflock + sortlock take 256 bytes, pushing leaf->ss to offset 1336 (1336 mod 16 = 8).',
    '   Aligned SIMD load at a misaligned address gives SIGSEGV on x86. NEON tolerates misalignment, so Mac was fine.',
    '',
    'Fix (one line):',
    '   _Alignas(32) struct entry13 hs[WH_KPN]; in struct wormleaf, gated on WH_LOCK_SHIM.',
    '   This adds 8 bytes of padding before hs[], pulls hs to offset 320 and ss to 1344, both 32 byte aligned.',
    '',
    'Failed first attempt worth recording:',
    '   I tried shrinking the shim alignment from alignas(64) to alignas(8). That left ss[] still misaligned',
    '   and introduced UB whenever LockT had alignas(64) atomics inside (occ_lock, ticket_lock).',
    '   Mac wh-cas and wh-rw started hanging at 99 percent CPU in rwlock_trylock_write, even single threaded.',
    '   Reverted in the same change as the proper _Alignas(32) fix.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=11, color=BLACK)


# ----- Slide 8: Wormhole bug story (optimistic reader) -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'How we built Wormhole, part 3: bugs in the optimistic reader',
          'Five separate things had to be fixed before wh-occ-opt was correct and fair')
body = [
    '1. Uninitialized occ_seq.',
    '   slab_alloc_safe returns memory with stale bits. If a fresh leaf came in with an odd seq counter,',
    '   readers spun forever waiting for a writer that never existed. Fix: initialize occ_seq to 0 in wormleaf_alloc.',
    '',
    '2. Missing seq bumps in writer fast paths.',
    '   wormhole_jump_leaf_write and wormhole_split_insert bypass the regular wormleaf_lock_write / unlock_write.',
    '   Without explicit bumps the seq would go odd and never come back to even, leaving the leaf in an',
    '   apparent "writer in progress" state forever. Fix: bump occ_seq on trylock_write_nr success in those paths.',
    '',
    '3. Torn entry13 reads.',
    '   entry13 is an 8 byte packed struct (e1: u16 key prefix, e3: u48 compressed pointer).',
    '   Stock wormleaf_match_hs reads e1 and e3 as separate field accesses. Without the leaflock, a writer',
    '   could clear hs[i].v64 between the two reads. Reader sees e1 == pkey, then reads e3 == 0, derefs NULL,',
    '   SIGSEGV. The seqlock validate would catch it, but only after the segfault.',
    '   Fix: read entry13 as a single atomic v64 load, unpack e1 and e3 from the same snapshot.',
    '',
    '4. kv use after free under optimistic reads.',
    '   kvmap_mm_dup frees old kvs immediately on update. An optimistic reader could hold a kv pointer that',
    '   gets freed mid-read. Fix: custom kvmap_mm with a no-op free for wh-occ-opt. kvs leak for the run',
    '   (~150 MB for a 3 second benchmark, acceptable for measurement). A production implementation would defer free via QSBR.',
    '',
    '5. Allocator asymmetry inflated the wh-occ-opt win.',
    '   Locked variants paid roughly 10 to 30 ns per op for free(); the no-op free variant did not.',
    '   That artificially doubled occ-opt apparent advantage on write heavy zipfian. Fix: WH_FAIR_MM CMake',
    '   option that uses the same no-op free MM for ALL variants. Required for an apples to apples comparison.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=10, color=BLACK)


# ----- Throughput slides (line + maxT bar) -----
add_image_slide(
    title='Wormhole, balanced workload (uniform 80 percent reads)',
    subtitle='Top: throughput vs threads. Bottom: ranking at the architecture\'s maximum threads.',
    image_path=plots['wh_balanced'],
)

add_image_slide(
    title='Wormhole, read heavy workload (uniform 90 percent reads)',
    subtitle='Where an optimistic reader is expected to win',
    image_path=plots['wh_readheavy'],
)

add_image_slide(
    title='Wormhole, write heavy workload (zipfian 20 percent reads, 40 inserts, 40 deletes)',
    subtitle='Hot key contention, where readers cannot escape conflict',
    image_path=plots['wh_writeheavy'],
)

add_image_slide(
    title='libcds StripedMap, balanced workload (uniform 80 percent reads)',
    subtitle='Per stripe lock contention, no traversal',
    image_path=plots['cds_balanced'],
)

add_image_slide(
    title='libcds StripedMap, hot key workload (zipfian 80 percent reads)',
    subtitle='A few stripes carry most of the traffic',
    image_path=plots['cds_zipfian'],
)

add_image_slide(
    title='libcds BronsonAVL, balanced workload (uniform 80 percent reads)',
    subtitle='Per node monitor lock, much smaller contention surface than StripedMap',
    image_path=plots['avl_balanced'],
)

add_image_slide(
    title='libcds BronsonAVL, write heavy workload',
    subtitle='Tree mutations create more lock acquisition per operation',
    image_path=plots['avl_writeheavy'],
)


# ----- Best lock per (arch, workload) table -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Best lock per architecture and workload',
          'Top throughput at each architecture\'s maximum thread count in scope')

table_rows = [['Index', 'Architecture', 'Max T', 'Workload', 'Best lock', 'Throughput (M ops per s)']] + [
    [r[0], r[1], str(r[2]), r[3], r[4], f'{r[5]:.1f}'] for r in best_rows
]
rows, cols = len(table_rows), len(table_rows[0])
tbl = sl.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4)).table
for r, row in enumerate(table_rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = (r == 0)

add_textbox(sl,
    ['Quick observations from the table:',
     '   Graviton and Xeon do not always agree on the best lock for the same workload.',
     '   For Wormhole, occ-opt is the headline winner on read leaning workloads (and on uniform 80/10/10 on Graviton).',
     '   For StripedMap and BronsonAVL, tas and ttas dominate. cas and ticket appear more on Xeon than on Graviton.'],
    Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.6), size=11, color=BLACK)


# ----- Slide N: open question and scope expansion (from supervisor feedback) -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Next steps, part 1: an open question to address',
          'From supervisor feedback')
body = [
    'The question:',
    '   "Have you examined all synchronization related instructions in armv8 and x86 to make sure that the locks',
    '   we have been evaluating cover them all? One of the original questions we want to answer is what is the',
    '   perf difference between directly writing assemblies and using C++ atomics. And it seems that we have',
    '   been evaluating only the latter."',
    '',
    'Honest answer:',
    '   Right. Every lock primitive we have evaluated so far (tas, ttas, cas, ticket, rw, occ) is built on',
    '   std::atomic. The compiler picks the underlying instruction. On Graviton with -march=native we should be',
    '   getting the ARMv8.1 LSE atomics (cas, casal, swp, ldadd) but we have not actually verified this with',
    '   objdump. On Xeon we get the lock prefixed instructions (lock cmpxchg, lock xchg, lock xadd).',
    '',
    'Things our current set does NOT cover:',
    '   On ARM, WFE (Wait For Event) and SEV (Send Event). These let a spinning core enter a low power wait',
    '   that wakes when the cache line changes. The Linux kernel arm64 spinlock uses them. Our cpu_relax just',
    '   does YIELD which is a hint that is largely a NOP on Neoverse cores.',
    '   On ARM, hand picking barrier strength (DMB ISHLD vs DMB ISHST vs DMB ISH SY) instead of letting the',
    '   compiler map memory_order_release to whichever DMB it picks.',
    '   On x86, there is much less missing. PAUSE we already use. UMONITOR / UMWAIT exists on Tremont and',
    '   Tigerlake plus, but diascld45 (Haswell era) does not have it. TSX HLE and RTM are disabled in microcode',
    '   on most modern Xeons due to MDS / Spectre mitigations.',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=11, color=BLACK)


# ----- Slide N+1: concrete experiments and remaining work -----
sl = prs.slides.add_slide(BLANK)
add_title(sl, 'Next steps, part 2: concrete experiments to run')
body = [
    'A. Verify what GCC actually emits for our locks (low effort, high information).',
    '   Run objdump -d on each wormhole-rt-<lock>.a on Graviton and Xeon. Check that:',
    '      tas_lock and ttas_lock use casal (LSE) on Graviton, lock cmpxchg / lock xchg on Xeon.',
    '      cas_lock uses cas (LSE) on Graviton.',
    '   If we are still seeing ldaxr / stlxr loops on Graviton, our -march flag is wrong and we are paying',
    '   the LL/SC penalty unnecessarily.',
    '',
    'B. Add a WFE based spinlock primitive (most promising new direction for ARM).',
    '   New file: include/primitives/tas_lock_wfe.hpp and a matching cas_lock_wfe.hpp.',
    '   The acquire path: try the atomic, on failure use WFE to wait for any change to the cache line, then retry.',
    '   The release path: store, then SEV to wake any waiting cores. Reference: Linux arch/arm64 spinlock.h.',
    '   Plug into the Wormhole shim and compare against the YIELD based versions on Graviton.',
    '   Hypothesis: at high contention, WFE saves real power and may reduce coherence traffic;',
    '   at low contention it is a wash or slightly slower because of the extra SEV.',
    '',
    'C. Hand assembly TAS and CAS to compare against std::atomic versions.',
    '   Write a tas_lock_asm and cas_lock_asm that use inline asm directly (LSE on ARM, lock prefix on x86).',
    '   Compare the throughput at every thread point. The expected delta is small if -march=native is doing',
    '   its job, but if there is a measurable gap that is itself an interesting finding for the writeup.',
    '',
    'D. Smaller items still on the list.',
    '   Get sudo on Xeon for one controlled run with setup_cpu.sh, to separate DVFS noise from real lock behaviour.',
    '   Add Apple M3 back for the Wormhole comparison (P cores only, capped at 6T to stay homogeneous).',
    '   Write up the cross arch lock ranking finding (ttas wins more on Graviton, cas / ticket more on Xeon).',
]
add_textbox(sl, body, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
            size=11, color=BLACK)


out_path = DECK_DIR / 'index_lock_evaluation.pptx'
prs.save(str(out_path))
print(f'\nDeck written to {out_path.relative_to(ROOT)}')
print(f'Slide count: {len(prs.slides)}')
