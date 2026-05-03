#!/usr/bin/env python3
"""Generate plots + PowerPoint deck summarizing the lockbench progress.

Usage: build_deck.py
Outputs:
  results/deck/figures/*.png
  results/deck/lockbench_progress.pptx
"""
from pathlib import Path
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results" / "avl_compare"
OUT_DIR = ROOT / "results" / "deck"
FIG_DIR = OUT_DIR / "figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR.mkdir(parents=True, exist_ok=True)

PALETTE = {"std": "#888888", "tas": "tab:blue", "ttas": "tab:green",
           "cas": "tab:orange", "ticket": "tab:red"}

# ---------------- Data ----------------
cds = pd.read_csv(RESULTS / "cds_striped.csv", sep=";", decimal=",")
avl = pd.read_csv(RESULTS / "cds_avl.csv",     sep=";", decimal=",")
cds["lk"] = cds["lock"].str.replace("cds-", "", regex=False)
avl["lk"] = avl["lock"].str.replace("avl-", "", regex=False)

def workload(row):
    rd, ins = int(row["read_pct"]), int(row["insert_pct"])
    if rd == 80 and ins == 10:  return f"{row['dist']} 80/10/10"
    if rd == 90 and ins == 5:   return "uniform 90/5/5 read-heavy"
    if rd == 20 and ins == 40:  return "zipfian 20/40/40 write-heavy"
    return f"{row['dist']} {rd}/{ins}/{100-rd-ins}"
for df in (cds, avl):
    df["workload"] = df.apply(workload, axis=1)

# ---------------- Plot helpers ----------------
def save_lock_grid(df, title, fname):
    workloads = sorted(df.workload.unique())
    fig, axes = plt.subplots(2, 2, figsize=(12, 8))
    for ax, wkl in zip(axes.flat, workloads):
        sub = df[df.workload == wkl]
        for lk in ["std", "tas", "ttas", "cas", "ticket"]:
            g = sub[sub.lk == lk].sort_values("threads")
            if len(g) == 0:
                continue
            ls = "--" if lk == "std" else "-"
            ax.plot(g.threads, g.ops_s / 1e6, marker="o",
                    color=PALETTE[lk], linestyle=ls, linewidth=2, label=lk)
        ax.set_title(wkl, fontsize=11)
        ax.set_xlabel("Threads")
        ax.set_ylabel("M ops/s")
        ax.legend(fontsize=9)
        ax.grid(True, alpha=0.3)
    fig.suptitle(title, fontsize=13, y=0.995)
    fig.tight_layout()
    fig.savefig(FIG_DIR / fname, dpi=130, bbox_inches="tight")
    plt.close(fig)

def save_speedup_table_plot(df, title, fname):
    """Bar plot at 8 threads showing speedup of each lock vs std."""
    workloads = sorted(df.workload.unique())
    fig, ax = plt.subplots(figsize=(11, 5))
    locks = ["std", "tas", "ttas", "cas", "ticket"]
    x_idx = list(range(len(workloads)))
    width = 0.16
    for i, lk in enumerate(locks):
        ys = []
        for wkl in workloads:
            g = df[(df.workload == wkl) & (df.lk == lk) & (df.threads == 8)]
            ys.append(g.ops_s.mean() / 1e6 if len(g) else 0)
        offset = (i - 2) * width
        ax.bar([x + offset for x in x_idx], ys, width=width,
               color=PALETTE[lk], label=lk, edgecolor="black", linewidth=0.4)
    ax.set_xticks(x_idx)
    ax.set_xticklabels([w.replace(" ", "\n", 1) for w in workloads], fontsize=9)
    ax.set_ylabel("M ops/s @ 8 threads")
    ax.set_title(title, fontsize=12)
    ax.legend(fontsize=9, ncol=5, loc="upper right")
    ax.grid(True, axis="y", alpha=0.3)
    fig.tight_layout()
    fig.savefig(FIG_DIR / fname, dpi=130, bbox_inches="tight")
    plt.close(fig)

print("Generating plots...")
save_lock_grid(cds, "libcds StripedMap — stripe-lock comparison", "cds_locks.png")
save_lock_grid(avl, "libcds BronsonAVLTreeMap — per-node lock comparison", "avl_locks.png")
save_speedup_table_plot(cds, "StripedMap — 8-thread throughput by lock", "cds_bar.png")
save_speedup_table_plot(avl, "BronsonAVL — 8-thread throughput by lock", "avl_bar.png")

# ---------------- Deck ----------------
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]
TITLE_LAYOUT = prs.slide_layouts[0]

def add_title(text, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    tx = s.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(2)).text_frame
    tx.word_wrap = True
    tx.text = text
    tx.paragraphs[0].runs[0].font.size = Pt(40)
    tx.paragraphs[0].runs[0].font.bold = True
    if subtitle:
        p = tx.add_paragraph()
        p.text = subtitle
        p.runs[0].font.size = Pt(20)
        p.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    return s

def add_section(title):
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1)).text_frame
    box.text = title
    box.paragraphs[0].runs[0].font.size = Pt(28)
    box.paragraphs[0].runs[0].font.bold = True
    return s

def add_bullets(title, bullets):
    s = add_section(title)
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5)).text_frame
    box.word_wrap = True
    box.text = ""  # first paragraph; will write bullets below
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            indent, txt = b
        else:
            indent, txt = 0, b
        p = box.paragraphs[0] if i == 0 else box.add_paragraph()
        p.text = txt
        p.level = indent
        for r in p.runs:
            r.font.size = Pt(20 if indent == 0 else 16)
    return s

def add_image(title, img_path, caption=None):
    s = add_section(title)
    pic = s.shapes.add_picture(str(img_path),
                               Inches(0.7), Inches(1.3),
                               width=Inches(12))
    # Cap height
    if pic.height > Inches(5.6):
        pic.height = Inches(5.6)
        pic.width  = Inches(5.6 * pic.width.inches / pic.height.inches)
        pic.left   = int((SW - pic.width) / 2)
    if caption:
        c = s.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.4)).text_frame
        c.text = caption
        c.paragraphs[0].runs[0].font.size = Pt(12)
        c.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        c.paragraphs[0].alignment = PP_ALIGN.CENTER
    return s

# ---- Title ----
add_title("Lockbench progress",
          "Trying different locks on libcds StripedMap, libcds AVL, and Wormhole")

# ---- What I'm doing ----
add_bullets("What I'm doing", [
    "I picked three concurrent indexes and tried plugging different lock primitives into each.",
    "Locks I'm using: std::mutex, TAS, TTAS, CAS, ticket.",
    "Indexes:",
    (1, "libcds StripedMap (hash map) — done"),
    (1, "libcds Bronson AVL tree — done"),
    (1, "Wormhole — still working on it"),
    "Same benchmark harness for all of them so the numbers are comparable.",
])

# ---- StripedMap ----
add_bullets("StripedMap setup", [
    "Vendored libcds (just the headers I need) into third_party/libcds/.",
    "The lock is a template parameter in libcds, so I just pass in my own lock types.",
    "Adapter: include/indexes/striped_map_index.hpp.",
    "Bench: bench/cds_bench.cpp, with --lock std|tas|ttas|cas|ticket.",
    "Wrote a correctness test (bench/cds_test.cpp) that compares against std::map.",
    "All 5 locks passed the correctness test.",
])

add_image("StripedMap throughput", FIG_DIR / "cds_locks.png",
          "Higher is better. std::mutex is the dashed grey line.")

add_image("StripedMap at 8 threads", FIG_DIR / "cds_bar.png")

add_bullets("StripedMap — what I noticed", [
    "TAS and TTAS are the fastest, especially on read-heavy.",
    "std::mutex is slower than the spinlocks here (Apple Silicon parks the thread).",
    "Ticket lock isn't great on uniform workloads.",
    "On the write-heavy workload everything looks similar — contention dominates.",
])

# ---- AVL ----
add_bullets("Bronson AVL setup", [
    "Same idea but with libcds's BronsonAVLTreeMap (a tree).",
    "Tree uses one lock per node, plumbed via cds::sync::injecting_monitor<Lock>.",
    "Needed to add a small static lib (cds-rt) for libcds's RCU runtime.",
    "Adapter: include/indexes/avl_tree_index.hpp. Bench: bench/cds_avl_bench.cpp.",
    "Correctness test passed for all 5 locks.",
])

add_image("Bronson AVL throughput", FIG_DIR / "avl_locks.png",
          "Slower than StripedMap overall because tree ops are O(log N).")

add_image("Bronson AVL at 8 threads", FIG_DIR / "avl_bar.png")

add_bullets("Bronson AVL — what I noticed", [
    "Ticket lock does really badly here.",
    "I think it's because Bronson's readers do a version-check thing and the ticket lock's FIFO ordering doesn't play well with that.",
    "TAS / TTAS / CAS are all close to each other.",
    "std::mutex is somewhere in the middle.",
])

# ---- Wormhole ----
add_bullets("Wormhole — what it is", [
    "Concurrent ordered index from Wu et al. (FAST '19).",
    "It's a mix of a trie + hash + linked list.",
    "Three lock sites in the code: leaflock (rwlock), metalock (rwlock), sortlock (small spinlock).",
    "Readers do a try-lock-with-spin then fall back to a version check if the lock isn't free.",
    "Also uses RCU (called QSBR) for reclaiming freed leaves — that's separate from the locks.",
])

add_bullets("Wormhole — plan", [
    "Vendored wormhole into third_party/wormhole/.",
    "Wrote a small shim that replaces wormhole's lock structs with one that wraps any of my lock types.",
    "Building one binary per lock (wh_bench_default, wh_bench_rw, wh_bench_tas, ...).",
    "7 variants total: default (wormhole's own lock) + rw, tas, ttas, cas, ticket, occ.",
    "Sweep script will run all 7 the same way as the StripedMap and AVL sweeps.",
])

add_bullets("Wormhole — things to watch out for", [
    "Ticket lock can't really do try-lock without messing up the queue, so I'm using an approximation.",
    "Under contention, ticket will probably look bad — that's a real finding, not a bug.",
    "OCC just runs as a write-only lock here (its optimistic-read protocol doesn't fit wormhole's API).",
    "Not touching wormhole's RCU code — that's separate from the locks.",
])

add_bullets("Wormhole — where I'm at", [
    "Done: vendored, wrote the shim, patched a few lines in lib.h / lib.c / wh.c, CMake builds all 7 variants.",
    "Had to fix one Apple Silicon asm thing in wormhole (semicolons vs newlines in .global directives).",
    "Now: running correctness tests for the 7 variants.",
    "Next: sweep script, plot, write up the caveats.",
])

# ---- Summary ----
add_bullets("Summary", [
    "StripedMap and Bronson AVL are done — locks plug in, correctness tests pass, plots show the differences.",
    "Same general pattern in both: TAS / TTAS are fastest, std::mutex pays parking cost, ticket struggles on the AVL.",
    "Wormhole is in progress — build works, correctness running, sweep next.",
])

out_file = OUT_DIR / "lockbench_progress.pptx"
prs.save(out_file)
print(f"Wrote {out_file}")
print(f"Figures in {FIG_DIR}")
